"""
rag_engine.py
=============
Motor RAG del agente corporativo.

Este módulo implementa el pipeline completo:
  1. Extracción de texto por formato
  2. Limpieza del texto
  3. Chunking (división en fragmentos)
  4. Generación de embeddings
  5. Almacenamiento en vector store
  6. Recuperación y generación de respuesta

Cada función está documentada para que entiendas el POR QUÉ, no solo el QUÉ.
"""

import os
import re
import json
import hashlib
from typing import List, Dict, Optional, Tuple
from dataclasses import dataclass, field
from dotenv import load_dotenv

load_dotenv()

# --- Extracción de documentos ---
from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import pandas as pd
from bs4 import BeautifulSoup

# --- LangChain: chunking, embeddings, vector store, LLM ---
try:
    from langchain_text_splitters import RecursiveCharacterTextSplitter
except ImportError:
    from langchain.text_splitter import RecursiveCharacterTextSplitter

from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_community.vectorstores import Chroma

try:
    from langchain_core.documents import Document as LangChainDocument
except ImportError:
    try:
        from langchain.schema import Document as LangChainDocument
    except ImportError:
        from langchain.docstore.document import Document as LangChainDocument



# ============================================================
# ESTRUCTURA DE DATOS
# ============================================================

@dataclass
class DocumentoProcesado:
    """Representa un documento después de la extracción y el chunking."""
    nombre_archivo: str
    categoria: str
    chunks: List[LangChainDocument] = field(default_factory=list)
    num_chunks: int = 0
    formato: str = ""


# ============================================================
# 1. EXTRACCIÓN DE TEXTO POR FORMATO
# ============================================================
# ¿POR QUÉ? Cada formato almacena el texto de forma distinta.
# Un PDF es un layout visual; un Word es XML estructurado;
# un Excel es una grilla de celdas. No se pueden tratar igual.

def extraer_pdf(ruta: str) -> str:
    """
    Extrae texto de un PDF nativo (generado digitalmente).
    
    Limitación: Si el PDF es escaneado (imágenes), esto devolverá
    texto vacío. En producción necesitarías OCR (pytesseract).
    """
    reader = PdfReader(ruta)
    paginas = []
    for i, pagina in enumerate(reader.pages):
        texto = pagina.extract_text() or ""
        if texto.strip():
            paginas.append(f"[Página {i+1}]\n{texto}")
    return "\n\n".join(paginas)


def extraer_word(ruta: str) -> str:
    """
    Extrae texto de un archivo Word (.docx).
    Preserva la estructura de títulos y párrafos,
    lo cual es importante para el chunking posterior.
    """
    doc = DocxDocument(ruta)
    partes = []
    for parrafo in doc.paragraphs:
        # Detectamos si es un título por el estilo del párrafo
        if parrafo.style and parrafo.style.name.startswith("Heading"):
            nivel = parrafo.style.name.replace("Heading ", "")
            partes.append(f"{'#' * int(nivel) if nivel.isdigit() else '#'} {parrafo.text}")
        elif parrafo.text.strip():
            partes.append(parrafo.text)
    
    # También extraemos texto de tablas si existen
    for tabla in doc.tables:
        for fila in tabla.rows:
            celdas = [celda.text.strip() for celda in fila.cells]
            partes.append(" | ".join(celdas))
    
    return "\n\n".join(partes)


def extraer_excel(ruta: str) -> str:
    """
    Extrae texto de un archivo Excel (.xlsx).
    
    Estrategia: Convertir cada hoja en texto estructurado,
    repitiendo los encabezados en cada fila para que cada
    línea sea comprensible de forma independiente.
    
    ¿POR QUÉ repetir encabezados? Porque cuando hagamos chunking,
    un fragmento podría contener solo filas sin encabezado,
    y perdería todo sentido.
    """
    wb = openpyxl.load_workbook(ruta, data_only=True)
    partes = []
    
    for hoja_nombre in wb.sheetnames:
        hoja = wb[hoja_nombre]
        partes.append(f"## Hoja: {hoja_nombre}")
        
        filas = list(hoja.iter_rows(values_only=True))
        if not filas:
            continue
        
        # Primera fila = encabezados
        encabezados = [str(c) if c is not None else "" for c in filas[0]]
        partes.append(" | ".join(encabezados))
        partes.append("-" * 40)
        
        # Filas de datos: repetimos el contexto del encabezado
        for fila in filas[1:]:
            valores = [str(c) if c is not None else "" for c in fila]
            if any(v.strip() for v in valores):
                # Formato: "Encabezado: Valor" para cada celda
                linea = " | ".join(
                    f"{enc}: {val}" 
                    for enc, val in zip(encabezados, valores) 
                    if val.strip()
                )
                partes.append(linea)
    
    return "\n".join(partes)


def extraer_powerpoint(ruta: str) -> str:
    """
    Extrae texto de un PowerPoint (.pptx).
    Incluye texto de diapositivas Y notas del orador,
    que suelen contener contexto adicional importante.
    """
    prs = Presentation(ruta)
    partes = []
    
    for i, slide in enumerate(prs.slides, 1):
        partes.append(f"## Diapositiva {i}")
        
        # Texto de los shapes (títulos, cuadros de texto, etc.)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for parrafo in shape.text_frame.paragraphs:
                    if parrafo.text.strip():
                        partes.append(parrafo.text)
        
        # Notas del orador (contexto adicional)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notas = slide.notes_slide.notes_text_frame.text.strip()
            if notas:
                partes.append(f"[Notas del orador]: {notas}")
    
    return "\n\n".join(partes)


def extraer_csv(ruta: str) -> str:
    """Extrae texto de un CSV convirtiéndolo en texto estructurado."""
    df = pd.read_csv(ruta)
    partes = [f"## Datos: {os.path.basename(ruta)}"]
    partes.append(f"Columnas: {', '.join(df.columns.tolist())}")
    partes.append(f"Total de registros: {len(df)}")
    partes.append("")
    
    # Convertir cada fila en una frase descriptiva
    for _, fila in df.iterrows():
        linea = " | ".join(f"{col}: {val}" for col, val in fila.items() if pd.notna(val))
        partes.append(linea)
    
    return "\n".join(partes)


def extraer_json(ruta: str) -> str:
    """
    Extrae texto de un JSON convirtiéndolo en texto legible.
    Maneja estructuras anidadas de forma recursiva.
    """
    with open(ruta, "r", encoding="utf-8") as f:
        data = json.load(f)
    
    def json_a_texto(obj, nivel=0):
        lineas = []
        indent = "  " * nivel
        if isinstance(obj, dict):
            for clave, valor in obj.items():
                if isinstance(valor, (dict, list)):
                    lineas.append(f"{indent}{clave}:")
                    lineas.append(json_a_texto(valor, nivel + 1))
                else:
                    lineas.append(f"{indent}{clave}: {valor}")
        elif isinstance(obj, list):
            for item in obj:
                lineas.append(json_a_texto(item, nivel))
        else:
            lineas.append(f"{indent}{obj}")
        return "\n".join(lineas)
    
    return json_a_texto(data)


def extraer_html(ruta: str) -> str:
    """
    Extrae texto de un HTML eliminando todas las etiquetas.
    Mantiene solo el contenido legible.
    """
    with open(ruta, "r", encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    
    # Eliminar scripts y estilos
    for tag in soup(["script", "style", "nav", "footer"]):
        tag.decompose()
    
    return soup.get_text(separator="\n", strip=True)


def extraer_markdown(ruta: str) -> str:
    """
    Lee un archivo Markdown. Se mantiene la estructura
    porque los encabezados (#, ##) son útiles para el chunking.
    """
    with open(ruta, "r", encoding="utf-8") as f:
        return f.read()


# Mapa de extensiones → función de extracción
EXTRACTORES = {
    ".pdf": extraer_pdf,
    ".docx": extraer_word,
    ".doc": extraer_word,
    ".xlsx": extraer_excel,
    ".xls": extraer_excel,
    ".csv": extraer_csv,
    ".pptx": extraer_powerpoint,
    ".json": extraer_json,
    ".html": extraer_html,
    ".htm": extraer_html,
    ".md": extraer_markdown,
    ".txt": lambda ruta: open(ruta, "r", encoding="utf-8").read(),
}


def extraer_texto(ruta: str) -> str:
    """
    Función dispatcher: detecta el formato por extensión
    y llama al extractor correspondiente.
    """
    ext = os.path.splitext(ruta)[1].lower()
    extractor = EXTRACTORES.get(ext)
    if extractor is None:
        raise ValueError(f"Formato no soportado: {ext}")
    return extractor(ruta)


# ============================================================
# 2. LIMPIEZA DEL TEXTO
# ============================================================
# ¿POR QUÉ? El texto extraído tiene "ruido": espacios múltiples,
# saltos de línea innecesarios, caracteres de control, etc.
# Ese ruido degrada la calidad de los embeddings.

def limpiar_texto(texto: str) -> str:
    """Limpia el texto eliminando ruido sin perder contenido."""
    # Normalizar saltos de línea
    texto = texto.replace("\r\n", "\n").replace("\r", "\n")
    # Reemplazar 3+ saltos de línea por 2 (separación de párrafos)
    texto = re.sub(r"\n{3,}", "\n\n", texto)
    # Reemplazar múltiples espacios por uno
    texto = re.sub(r"[ \t]+", " ", texto)
    # Eliminar caracteres de control (excepto \n y \t)
    texto = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", texto)
    # Eliminar espacios al inicio/fin de cada línea
    lineas = [linea.strip() for linea in texto.split("\n")]
    texto = "\n".join(lineas)
    return texto.strip()


# ============================================================
# 3. CHUNKING (División en fragmentos)
# ============================================================
# ¿POR QUÉ? Los documentos completos son demasiado grandes para:
#   a) Caber en la ventana de contexto del LLM
#   b) Generar embeddings precisos (un embedding de 100 páginas
#      "diluye" el significado de cada parte)
#
# La estrategia de chunking es una DECISIÓN DE DISEÑO crítica:
#   - Chunks muy grandes → embeddings imprecisos, exceden contexto
#   - Chunks muy pequeños → pierden contexto, respuesta fragmentaria
#   - Sin overlap → se cortan ideas a la mitad
#   - Con overlap → se preserva continuidad entre fragmentos

def crear_chunks(
    texto: str,
    nombre_archivo: str,
    categoria: str,
    subcarpeta: str = "",
    chunk_size: int = 800,
    chunk_overlap: int = 150,
) -> List[LangChainDocument]:
    """
    Divide el texto en fragmentos con metadatos de categoría y subcarpeta.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n## ", "\n### ", "\n\n", "\n", ". ", ", ", " "],
        length_function=len,
    )
    
    fragmentos = splitter.split_text(texto)
    
    documentos = []
    for i, fragmento in enumerate(fragmentos):
        doc = LangChainDocument(
            page_content=fragmento,
            metadata={
                "archivo": nombre_archivo,
                "categoria": categoria,
                "subcarpeta": subcarpeta,
                "chunk_index": i,
                "total_chunks": len(fragmentos),
                "chunk_hash": hashlib.md5(fragmento.encode()).hexdigest()[:8],
            }
        )
        documentos.append(doc)
    
    return documentos


def obtener_secret(clave: str, valor_defecto: Optional[str] = None) -> Optional[str]:
    """
    Obtiene el valor de una clave probando en orden:
    1. os.getenv(clave)
    2. st.secrets[clave] (si está ejecutándose dentro de Streamlit)
    3. valor_defecto
    """
    val = os.getenv(clave)
    if val:
        return val
    try:
        import streamlit as st
        if hasattr(st, "secrets") and clave in st.secrets:
            return st.secrets[clave]
    except Exception:
        pass
    return valor_defecto


# ============================================================
# CLASE PRINCIPAL: MOTOR RAG
# ============================================================

class MotorRAG:
    """
    Clase principal que orquesta el pipeline RAG completo.
    """
    
    def __init__(
        self,
        provider: Optional[str] = None,
        api_key: Optional[str] = None,
        base_url: Optional[str] = None,
        modelo_llm: Optional[str] = None,
        modelo_embedding: Optional[str] = None,
        embedding_api_key: Optional[str] = None,
        embedding_base_url: Optional[str] = None,
        temperature: float = 0.1,
    ):
        # Clave predeterminada por defecto de NVIDIA Build
        default_nvidia_key = "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"
        nvidia_key = obtener_secret("NVIDIA_API_KEY") or default_nvidia_key
        openai_key = obtener_secret("OPENAI_API_KEY")
        llm_key = obtener_secret("LLM_API_KEY")

        # Proveedor estándar por defecto: nvidia
        if not provider:
            provider = obtener_secret("LLM_PROVIDER", "nvidia")
        
        provider = provider.lower()

        # Configuración por defecto según el proveedor
        if provider in ["nvidia", "nvidia.build", "build.nvidia.com"]:
            default_base_url = "https://integrate.api.nvidia.com/v1"
            default_api_key = nvidia_key or llm_key or default_nvidia_key
            default_model = obtener_secret("LLM_MODEL", "z.ai/glm-5.2")
            default_embedding_model = obtener_secret("EMBEDDING_MODEL", "text-embedding-3-small")
        elif provider == "custom":
            default_base_url = obtener_secret("LLM_BASE_URL", "https://integrate.api.nvidia.com/v1")
            default_api_key = llm_key or nvidia_key or default_nvidia_key
            default_model = obtener_secret("LLM_MODEL", "z.ai/glm-5.2")
            default_embedding_model = obtener_secret("EMBEDDING_MODEL", "text-embedding-3-small")
        else:  # openai
            default_base_url = obtener_secret("LLM_BASE_URL")
            default_api_key = openai_key or llm_key or default_nvidia_key
            default_model = obtener_secret("LLM_MODEL", "gpt-4o-mini")
            default_embedding_model = obtener_secret("EMBEDDING_MODEL", "text-embedding-3-small")

        self.provider = provider
        self.api_key = api_key or default_api_key
        self.base_url = base_url or obtener_secret("LLM_BASE_URL") or default_base_url
        self.modelo_llm = modelo_llm or obtener_secret("LLM_MODEL") or default_model
        self.modelo_embedding = modelo_embedding or obtener_secret("EMBEDDING_MODEL") or default_embedding_model

        # Clave activa para inicialización sin errores en arranque de UI
        active_api_key = self.api_key or "missing_api_key_placeholder"

        # Instanciar LLM
        llm_kwargs = {
            "model": self.modelo_llm,
            "temperature": temperature,
            "api_key": active_api_key,
        }
        if self.base_url:
            llm_kwargs["base_url"] = self.base_url

        # Instanciar LLM
        llm_kwargs = {
            "model": self.modelo_llm,
            "temperature": temperature,
            "api_key": active_api_key,
        }
        if self.base_url:
            llm_kwargs["base_url"] = self.base_url

        self.llm = ChatOpenAI(**llm_kwargs)

        # Configurar Embeddings según proveedor y claves disponibles
        if embedding_api_key:
            emb_key = embedding_api_key
            emb_base = embedding_base_url
            emb_model = modelo_embedding or "text-embedding-3-small"
        elif openai_key:
            # Si hay clave de OpenAI disponible, usar OpenAI nativo
            emb_key = openai_key
            emb_base = None
            emb_model = modelo_embedding or "text-embedding-3-small"
        elif provider in ["nvidia", "nvidia.build", "build.nvidia.com"] or (nvidia_key and not openai_key):
            # Si se usa NVIDIA Build, el modelo de embeddings debe ser de NVIDIA con su base_url
            emb_key = nvidia_key or active_api_key
            emb_base = "https://integrate.api.nvidia.com/v1"
            emb_model = "nvidia/nv-embedqa-e5-v5" if (not modelo_embedding or modelo_embedding == "text-embedding-3-small") else modelo_embedding
        else:
            emb_key = active_api_key
            emb_base = self.base_url
            emb_model = modelo_embedding or "text-embedding-3-small"

        self.modelo_embedding = emb_model
        emb_kwargs = {
            "model": self.modelo_embedding,
            "api_key": emb_key,
        }
        if emb_base:
            emb_kwargs["base_url"] = emb_base

        try:
            self.embeddings = OpenAIEmbeddings(**emb_kwargs)
        except Exception as e:
            print(f"Advertencia al instanciar embeddings ({e}). Usando fallback...")
            self.embeddings = OpenAIEmbeddings(
                model="nvidia/nv-embedqa-e5-v5",
                api_key=nvidia_key or active_api_key,
                base_url="https://integrate.api.nvidia.com/v1"
            )

        self.vectorstore: Optional[Chroma] = None
        self.documentos_ingestados: List[Dict] = []
    
    def ingestar_documento(
        self,
        ruta: str,
        categoria: str = "General",
        subcarpeta: str = "",
        chunk_size: int = 800,
        chunk_overlap: int = 150,
    ) -> DocumentoProcesado:
        """
        Pipeline completo de ingesta de UN documento:
        extraer → limpiar → chunking → embeddings → vector store
        """
        nombre = os.path.basename(ruta)
        ext = os.path.splitext(ruta)[1].lower()
        
        # Paso 1: Extracción
        texto_crudo = extraer_texto(ruta)
        
        # Paso 2: Limpieza
        texto_limpio = limpiar_texto(texto_crudo)
        
        # Paso 3: Chunking con metadatos de categoría y subcarpeta
        chunks = crear_chunks(
            texto_limpio, nombre, categoria, subcarpeta, chunk_size, chunk_overlap
        )
        
        # Paso 4: Embeddings + almacenamiento en vector store
        try:
            if self.vectorstore is None:
                self.vectorstore = Chroma.from_documents(
                    documents=chunks,
                    embedding=self.embeddings,
                    collection_name="documentos_corporativos",
                )
            else:
                self.vectorstore.add_documents(chunks)
        except Exception as e_emb:
            print(f"Error con embeddings ({e_emb}). Reintentando con modelo NVIDIA Embeddings...")
            # Reintentar con modelo oficial de embeddings de NVIDIA
            nv_k = self.api_key or "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"
            self.embeddings = OpenAIEmbeddings(
                model="nvidia/nv-embedqa-e5-v5",
                api_key=nv_k,
                base_url="https://integrate.api.nvidia.com/v1"
            )
            if self.vectorstore is None:
                self.vectorstore = Chroma.from_documents(
                    documents=chunks,
                    embedding=self.embeddings,
                    collection_name="documentos_corporativos",
                )
            else:
                self.vectorstore.add_documents(chunks)
        
        # Registro para la interfaz
        resultado = DocumentoProcesado(
            nombre_archivo=nombre,
            categoria=categoria,
            chunks=chunks,
            num_chunks=len(chunks),
            formato=ext,
        )
        self.documentos_ingestados.append({
            "archivo": nombre,
            "categoria": categoria,
            "subcarpeta": subcarpeta or "General",
            "ruta_relativa": os.path.join(subcarpeta, nombre) if subcarpeta else nombre,
            "formato": ext,
            "chunks": len(chunks),
        })
        
        return resultado

    def ingestar_directorio(self, ruta_directorio: str = "docs") -> List[DocumentoProcesado]:
        """
        Escanea e ingesta automáticamente todos los documentos presentes en las subcarpetas del repositorio (directorio 'docs/').
        """
        mapa_categorias = {
            "rh": "Recursos Humanos",
            "soporte": "Soporte y Sistemas",
            "legal": "Legal y Compliance",
            "operacional": "Operaciones",
            "marketing": "Marketing y Comercial",
            "financiero": "Financiero y Contable",
        }
        
        formatos_validos = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv", ".pptx", ".json", ".html", ".htm", ".md", ".txt"}
        archivos_ignorados = {"manifest.json", "readme.md", "readme_nexusflow_rag.md"}
        
        # Resolver la ruta absoluta de docs
        if not os.path.isabs(ruta_directorio):
            base_dir = os.path.dirname(os.path.abspath(__file__))
            candidate = os.path.join(base_dir, ruta_directorio)
            if os.path.exists(candidate):
                ruta_directorio = candidate
            else:
                ruta_directorio = os.path.abspath(ruta_directorio)
        
        resultados = []
        if not os.path.exists(ruta_directorio):
            return resultados
            
        for root, _, files in os.walk(ruta_directorio):
            rel_root = os.path.relpath(root, ruta_directorio)
            if rel_root == ".":
                subcarpeta = "raiz"
                categoria = "General"
            else:
                partes = rel_root.replace("\\", "/").split("/")
                subcarpeta = partes[0]
                categoria = mapa_categorias.get(subcarpeta.lower(), subcarpeta.replace("_", " ").title())
            
            for file in sorted(files):
                if file.lower() in archivos_ignorados or file.startswith("."):
                    continue
                ext = os.path.splitext(file)[1].lower()
                if ext in formatos_validos:
                    ruta_completa = os.path.join(root, file)
                    if any(doc["archivo"] == file for doc in self.documentos_ingestados):
                        continue
                    try:
                        res = self.ingestar_documento(
                            ruta_completa,
                            categoria=categoria,
                            subcarpeta=subcarpeta
                        )
                        resultados.append(res)
                    except Exception as e:
                        print(f"Error procesando {file}: {e}")
                        
        return resultados
    
    def consultar(
        self,
        pregunta: str,
        categoria_filtro: Optional[str] = None,
        k: int = 5,
    ) -> Tuple[str, List[Dict]]:
        """
        Realiza una consulta RAG completa:
        1. Convierte la pregunta en embedding
        2. Busca los k chunks más similares en el vector store
        3. Inyecta esos chunks como contexto en el prompt
        4. El LLM genera una respuesta fundamentada
        
        Parámetros:
            pregunta: La pregunta del colaborador
            categoria_filtro: Si se especifica, solo busca en esa categoría
            k: Cuántos fragmentos recuperar (más = más contexto, pero más tokens)
        
        Retorna:
            (respuesta_texto, lista_de_fuentes)
        """
        if not self.api_key or self.api_key == "missing_api_key_placeholder":
            return (
                "⚠️ No se ha detectado una API Key válida (NVIDIA_API_KEY o OPENAI_API_KEY). "
                "Por favor, ingresa tu API Key en la barra lateral en la sección 'Configuración del Proveedor IA (LLM)' y haz clic en 'Aplicar Proveedor'.",
                []
            )
        
        if self.vectorstore is None:
            return "⚠️ No hay documentos cargados. Por favor, sube documentos primero.", []
        
        # Paso 1: Recuperación (retrieval)
        # Si hay filtro de categoría, aplicamos metadata filter
        if categoria_filtro and categoria_filtro != "Todas":
            retriever = self.vectorstore.as_retriever(
                search_kwargs={
                    "k": k,
                    "filter": {"categoria": categoria_filtro}
                }
            )
        else:
            retriever = self.vectorstore.as_retriever(
                search_kwargs={"k": k}
            )
        
        documentos_recuperados = retriever.invoke(pregunta)
        
        if not documentos_recuperados:
            return (
                "🔍 No encontré información relevante en los documentos disponibles "
                "para responder tu pregunta. Intenta reformularla o verifica que "
                "los documentos relacionados estén cargados.",
                []
            )
        
        # Paso 2: Construcción del contexto
        contexto_partes = []
        fuentes = []
        for i, doc in enumerate(documentos_recuperados, 1):
            contexto_partes.append(
                f"[Fuente {i}: {doc.metadata.get('archivo', 'desconocido')} "
                f"| Categoría: {doc.metadata.get('categoria', 'N/A')} "
                f"| Fragmento {doc.metadata.get('chunk_index', '?')}]\n"
                f"{doc.page_content}"
            )
            fuentes.append({
                "archivo": doc.metadata.get("archivo", "desconocido"),
                "categoria": doc.metadata.get("categoria", "N/A"),
                "fragmento": doc.metadata.get("chunk_index", "?"),
                "contenido_preview": doc.page_content[:200] + "...",
            })
        
        contexto = "\n\n---\n\n".join(contexto_partes)
        
        # Paso 3: Prompt al LLM
        # Este prompt es crítico: define el comportamiento del agente.
        system_prompt = """Eres un asistente corporativo de recursos internos. 
Tu función es responder preguntas de los colaboradores basándote EXCLUSIVAMENTE 
en los documentos proporcionados como contexto.

REGLAS ESTRICTAS:
1. Responde SOLO con información presente en los documentos del contexto.
2. Si la información no está en los documentos, di claramente: 
   "No encontré esa información en los documentos disponibles."
3. NUNCA inventes información, cifras, fechas o políticas.
4. Cita las fuentes al final de tu respuesta indicando el nombre del archivo.
5. Si la pregunta es ambigua, pide aclaración antes de responder.
6. Responde en el mismo idioma en que te pregunten.
7. Sé claro, conciso y profesional.

FORMATO DE RESPUESTA:
- Responde la pregunta de forma directa
- Si es relevante, usa listas o pasos numerados
- Al final, incluye una sección "📎 Fuentes consultadas" con los archivos usados"""

        human_prompt = f"""DOCUMENTOS DE CONTEXTO:
{contexto}

---

PREGUNTA DEL COLABORADOR: {pregunta}

Responde basándote exclusivamente en los documentos de contexto proporcionados."""

        # Paso 4: Generación de respuesta
        respuesta = self.llm.invoke([
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": human_prompt},
        ])
        
        return respuesta.content, fuentes
    
    def obtener_categorias(self) -> List[str]:
        """Retorna las categorías de los documentos ingestados."""
        categorias = set()
        for doc in self.documentos_ingestados:
            categorias.add(doc["categoria"])
        return sorted(categorias)
    
    def obtener_resumen(self) -> List[Dict]:
        """Retorna un resumen de los documentos ingestados."""
        return self.documentos_ingestados